"""
Enhanced File Processing Utilities
Supports multiple file types for better document processing capabilities.
"""

import logging
import os
import tempfile
import io
import json
from typing import Dict, List, Any, Optional, Tuple
from pathlib import Path

# Additional imports for file processing
try:
    import docx
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import markdown
    from bs4 import BeautifulSoup
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

class EnhancedFileProcessor:
    """Enhanced file processing with support for multiple formats."""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        
        # Supported file extensions and their processors
        self.processors = {
            # Document formats
            '.txt': self._process_text_file,
            '.md': self._process_markdown_file,
            '.json': self._process_json_file,
            '.csv': self._process_csv_file,
            '.xlsx': self._process_excel_file,
            '.xls': self._process_excel_file,
            
            # Code files
            '.py': self._process_code_file,
            '.js': self._process_code_file,
            '.html': self._process_html_file,
            '.css': self._process_code_file,
            '.xml': self._process_xml_file,
            '.yaml': self._process_yaml_file,
            '.yml': self._process_yaml_file,
            
            # Log files
            '.log': self._process_log_file,
        }
        
        # Add conditional processors
        if DOCX_AVAILABLE:
            self.processors['.docx'] = self._process_docx_file
        
        if PPTX_AVAILABLE:
            self.processors['.pptx'] = self._process_pptx_file
    
    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        return list(self.processors.keys())
    
    def is_supported(self, filename: str) -> bool:
        """Check if file type is supported."""
        ext = Path(filename).suffix.lower()
        return ext in self.processors
    
    async def process_file(self, file_path: str, original_name: str = None) -> Dict[str, Any]:
        """
        Process a file and extract its content and metadata.
        
        Args:
            file_path (str): Path to the file
            original_name (str): Original filename
            
        Returns:
            Dict[str, Any]: Processing result
        """
        try:
            if not os.path.exists(file_path):
                return {
                    "success": False,
                    "error": "File not found",
                    "content": None,
                    "metadata": {}
                }
            
            # Determine file extension
            filename = original_name or os.path.basename(file_path)
            ext = Path(filename).suffix.lower()
            
            if ext not in self.processors:
                return {
                    "success": False,
                    "error": f"Unsupported file type: {ext}",
                    "content": None,
                    "metadata": {"filename": filename, "extension": ext}
                }
            
            # Get file stats
            file_stats = os.stat(file_path)
            metadata = {
                "filename": filename,
                "extension": ext,
                "size_bytes": file_stats.st_size,
                "size_human": self._format_file_size(file_stats.st_size),
                "processor": self.processors[ext].__name__
            }
            
            # Process the file
            processor = self.processors[ext]
            content, extra_metadata = await processor(file_path, filename)
            
            metadata.update(extra_metadata)
            
            return {
                "success": True,
                "error": None,
                "content": content,
                "metadata": metadata
            }
            
        except Exception as e:
            self.logger.error(f"Error processing file {file_path}: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "content": None,
                "metadata": {"filename": filename if 'filename' in locals() else "unknown"}
            }
    
    async def _process_text_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            metadata = {
                "lines": len(content.splitlines()),
                "characters": len(content),
                "words": len(content.split()),
                "type": "text"
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"Error reading text file: {str(e)}")
    
    async def _process_markdown_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process Markdown files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                md_content = f.read()
            
            # Convert to HTML if markdown library is available
            html_content = None
            if MARKDOWN_AVAILABLE:
                try:
                    html_content = markdown.markdown(md_content, extensions=['tables', 'fenced_code'])
                    if BeautifulSoup:
                        soup = BeautifulSoup(html_content, 'html.parser')
                        headers = [h.get_text() for h in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])]
                        links = [a.get('href') for a in soup.find_all('a', href=True)]
                    else:
                        headers = []
                        links = []
                except Exception:
                    headers = []
                    links = []
            else:
                headers = []
                links = []
            
            metadata = {
                "lines": len(md_content.splitlines()),
                "characters": len(md_content),
                "words": len(md_content.split()),
                "type": "markdown",
                "headers": headers[:10],  # Limit to first 10 headers
                "links": links[:20],  # Limit to first 20 links
                "has_html": html_content is not None
            }
            
            # Return both markdown and HTML if available
            if html_content:
                content = f"# Markdown Content:\n{md_content}\n\n# HTML Preview:\n{html_content}"
            else:
                content = md_content
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"Error reading markdown file: {str(e)}")
    
    async def _process_json_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process JSON files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                raw_content = f.read()
                json_data = json.loads(raw_content)
            
            # Format JSON for better readability
            formatted_json = json.dumps(json_data, indent=2, ensure_ascii=False)
            
            metadata = {
                "type": "json",
                "structure": self._analyze_json_structure(json_data),
                "size_formatted": len(formatted_json),
                "is_valid": True
            }
            
            return formatted_json, metadata
            
        except json.JSONDecodeError as e:
            # Return raw content if JSON is invalid
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                raw_content = f.read()
            
            metadata = {
                "type": "json",
                "is_valid": False,
                "error": str(e)
            }
            
            return raw_content, metadata
    
    async def _process_csv_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process CSV files."""
        try:
            import pandas as pd
            
            # Read CSV with pandas for better analysis
            df = pd.read_csv(file_path, encoding='utf-8', errors='ignore', nrows=100)  # Limit to first 100 rows
            
            # Generate summary
            content_lines = [
                f"CSV File Analysis: {filename}",
                f"Shape: {df.shape[0]} rows, {df.shape[1]} columns",
                "",
                "Columns:",
                *[f"  - {col} ({str(df[col].dtype)})" for col in df.columns],
                "",
                "First 10 rows:",
                df.head(10).to_string(),
                "",
                "Data Info:",
                str(df.describe(include='all').to_string()) if not df.empty else "No data to describe"
            ]
            
            metadata = {
                "type": "csv",
                "rows": len(df),
                "columns": len(df.columns),
                "column_names": list(df.columns),
                "data_types": {col: str(dtype) for col, dtype in df.dtypes.items()},
                "memory_usage": df.memory_usage(deep=True).sum(),
                "has_null_values": df.isnull().any().any()
            }
            
            return "\n".join(content_lines), metadata
            
        except Exception as e:
            # Fallback to simple text processing
            return await self._process_text_file(file_path, filename)
    
    async def _process_excel_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process Excel files."""
        try:
            import pandas as pd
            
            # Read Excel file
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names
            
            content_lines = [f"Excel File Analysis: {filename}", ""]
            
            metadata = {
                "type": "excel",
                "sheets": sheet_names,
                "sheet_count": len(sheet_names)
            }
            
            # Process each sheet (limit to first 3 sheets)
            for i, sheet_name in enumerate(sheet_names[:3]):
                df = pd.read_excel(file_path, sheet_name=sheet_name, nrows=50)
                
                content_lines.extend([
                    f"Sheet: {sheet_name}",
                    f"Shape: {df.shape[0]} rows, {df.shape[1]} columns",
                    f"Columns: {', '.join(df.columns)}",
                    "",
                    "First 5 rows:",
                    df.head(5).to_string(),
                    "",
                    "---",
                    ""
                ])
                
                metadata[f"sheet_{i}"] = {
                    "name": sheet_name,
                    "rows": len(df),
                    "columns": len(df.columns),
                    "column_names": list(df.columns)
                }
            
            return "\n".join(content_lines), metadata
            
        except Exception as e:
            raise Exception(f"Error processing Excel file: {str(e)}")
    
    async def _process_docx_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process Word documents."""
        if not DOCX_AVAILABLE:
            raise Exception("python-docx library not available")
        
        try:
            doc = Document(file_path)
            
            # Extract text content
            content_lines = []
            paragraph_count = 0
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content_lines.append(paragraph.text)
                    paragraph_count += 1
            
            # Extract tables
            table_count = len(doc.tables)
            if table_count > 0:
                content_lines.append("\n--- TABLES ---")
                for i, table in enumerate(doc.tables[:3]):  # Limit to first 3 tables
                    content_lines.append(f"\nTable {i+1}:")
                    for row in table.rows[:5]:  # First 5 rows
                        row_text = " | ".join([cell.text for cell in row.cells])
                        content_lines.append(row_text)
            
            content = "\n".join(content_lines)
            
            metadata = {
                "type": "docx",
                "paragraphs": paragraph_count,
                "tables": table_count,
                "words": len(content.split()),
                "characters": len(content)
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"Error processing Word document: {str(e)}")
    
    async def _process_pptx_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process PowerPoint presentations."""
        if not PPTX_AVAILABLE:
            raise Exception("python-pptx library not available")
        
        try:
            prs = Presentation(file_path)
            
            content_lines = [f"PowerPoint Presentation: {filename}", ""]
            slide_count = 0
            
            for i, slide in enumerate(prs.slides):
                slide_count += 1
                content_lines.append(f"--- Slide {i+1} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_lines.append(shape.text)
                
                content_lines.append("")
            
            content = "\n".join(content_lines)
            
            metadata = {
                "type": "pptx",
                "slides": slide_count,
                "words": len(content.split()),
                "characters": len(content)
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"Error processing PowerPoint file: {str(e)}")
    
    async def _process_code_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process code files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # Basic code analysis
            lines = content.splitlines()
            non_empty_lines = [line for line in lines if line.strip()]
            comment_lines = [line for line in lines if line.strip().startswith(('#', '//', '/*', '*', '<!--'))]
            
            metadata = {
                "type": "code",
                "language": Path(filename).suffix[1:],  # Remove the dot
                "total_lines": len(lines),
                "code_lines": len(non_empty_lines),
                "comment_lines": len(comment_lines),
                "characters": len(content)
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"Error reading code file: {str(e)}")
    
    async def _process_html_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process HTML files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            # Extract text content if BeautifulSoup is available
            text_content = html_content
            if BeautifulSoup:
                try:
                    soup = BeautifulSoup(html_content, 'html.parser')
                    text_content = soup.get_text()
                    
                    # Extract metadata
                    title = soup.find('title')
                    meta_tags = soup.find_all('meta')
                    links = soup.find_all('a', href=True)
                    
                    metadata = {
                        "type": "html",
                        "title": title.text if title else "No title",
                        "meta_count": len(meta_tags),
                        "link_count": len(links),
                        "characters": len(text_content),
                        "has_parsed_content": True
                    }
                except Exception:
                    metadata = {"type": "html", "has_parsed_content": False}
            else:
                metadata = {"type": "html", "has_parsed_content": False}
            
            # Return both HTML and text content
            content = f"HTML Content:\n{html_content}\n\nExtracted Text:\n{text_content}"
            return content, metadata
            
        except Exception as e:
            raise Exception(f"Error reading HTML file: {str(e)}")
    
    async def _process_xml_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process XML files."""
        return await self._process_text_file(file_path, filename)
    
    async def _process_yaml_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process YAML files."""
        try:
            import yaml
            
            with open(file_path, 'r', encoding='utf-8') as f:
                yaml_content = f.read()
                yaml_data = yaml.safe_load(yaml_content)
            
            # Format YAML for better readability
            formatted_yaml = yaml.dump(yaml_data, default_flow_style=False, allow_unicode=True)
            
            metadata = {
                "type": "yaml",
                "structure": self._analyze_json_structure(yaml_data),  # Same analysis as JSON
                "is_valid": True
            }
            
            return formatted_yaml, metadata
            
        except Exception:
            # Fallback to text processing
            return await self._process_text_file(file_path, filename)
    
    async def _process_log_file(self, file_path: str, filename: str) -> Tuple[str, Dict]:
        """Process log files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                lines = f.readlines()
            
            # Analyze log patterns (basic)
            error_lines = [line for line in lines if any(keyword in line.lower() for keyword in ['error', 'exception', 'fail'])]
            warning_lines = [line for line in lines if any(keyword in line.lower() for keyword in ['warn', 'warning'])]
            
            # Show last 50 lines for recent activity
            recent_lines = lines[-50:] if len(lines) > 50 else lines
            content = "".join(recent_lines)
            
            metadata = {
                "type": "log",
                "total_lines": len(lines),
                "error_lines": len(error_lines),
                "warning_lines": len(warning_lines),
                "showing_recent": len(recent_lines)
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"Error reading log file: {str(e)}")
    
    def _analyze_json_structure(self, data: Any, max_depth: int = 3) -> Dict[str, Any]:
        """Analyze JSON structure."""
        if max_depth <= 0:
            return {"type": type(data).__name__, "truncated": True}
        
        if isinstance(data, dict):
            return {
                "type": "object",
                "keys": list(data.keys())[:10],  # First 10 keys
                "key_count": len(data.keys()),
                "sample_values": {k: self._analyze_json_structure(v, max_depth-1) 
                                 for k, v in list(data.items())[:3]}  # First 3 key-value pairs
            }
        elif isinstance(data, list):
            return {
                "type": "array",
                "length": len(data),
                "sample_items": [self._analyze_json_structure(item, max_depth-1) 
                               for item in data[:3]] if data else []
            }
        else:
            return {"type": type(data).__name__, "value": str(data)[:100]}
    
    def _format_file_size(self, size_bytes: int) -> str:
        """Format file size in human-readable format."""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024.0:
                return f"{size_bytes:.1f} {unit}"
            size_bytes /= 1024.0
        return f"{size_bytes:.1f} TB"

# Global instance
enhanced_file_processor = EnhancedFileProcessor()